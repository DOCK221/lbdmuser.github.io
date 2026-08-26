#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Couleurs du thème
    COLOR_PRIMARY = RGBColor(41, 98, 255)  # Bleu professionnel
    COLOR_SECONDARY = RGBColor(255, 87, 34)  # Orange dynamique
    COLOR_TEXT = RGBColor(51, 51, 51)  # Gris foncé
    COLOR_BG = RGBColor(245, 245, 245)  # Gris clair
    
    def add_title_slide(title, subtitle=""):
        """Crée une slide de titre"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Fond de couleur
        background = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = COLOR_PRIMARY
        background.line.color.rgb = COLOR_PRIMARY
        
        # Titre
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Sous-titre
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(8), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items):
        """Crée une slide de contenu"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Fond blanc
        background = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background.line.color.rgb = RGBColor(255, 255, 255)
        
        # Barre de titre colorée
        title_bar = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = COLOR_PRIMARY
        title_bar.line.color.rgb = COLOR_PRIMARY
        
        # Titre
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Contenu
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(12)
            p.level = 0
        
        return slide
    
    def add_highlight_slide(title, highlight_text, details=None):
        """Crée une slide avec mise en avant"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Fond blanc
        background = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background.line.color.rgb = RGBColor(255, 255, 255)
        
        # Titre
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
        
        # Encadré de mise en avant
        highlight_box = slide.shapes.add_shape(
            1, Inches(1.5), Inches(1.8), Inches(7), Inches(2)
        )
        highlight_box.fill.solid()
        highlight_box.fill.fore_color.rgb = COLOR_SECONDARY
        highlight_box.line.color.rgb = COLOR_SECONDARY
        
        # Texte mis en avant
        highlight_text_box = slide.shapes.add_textbox(
            Inches(1.8), Inches(2.2), Inches(6.4), Inches(1.2)
        )
        highlight_frame = highlight_text_box.text_frame
        highlight_frame.text = highlight_text
        highlight_frame.paragraphs[0].font.size = Pt(24)
        highlight_frame.paragraphs[0].font.bold = True
        highlight_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        highlight_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Détails si fournis
        if details:
            details_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(8), Inches(2.5)
            )
            details_frame = details_box.text_frame
            for i, detail in enumerate(details):
                if i > 0:
                    details_frame.add_paragraph()
                p = details_frame.paragraphs[i]
                p.text = detail
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_TEXT
                p.space_after = Pt(10)
        
        return slide
    
    # ==========================================
    # SLIDE 1 : PAGE DE GARDE
    # ==========================================
    slide1 = add_title_slide(
        "Retour d'Expérience Professionnelle",
        "Community Manager - Buroxia"
    )
    
    # Ajout des infos en bas
    info_box = slide1.shapes.add_textbox(
        Inches(2), Inches(5.5), Inches(6), Inches(1.5)
    )
    info_frame = info_box.text_frame
    info_text = "Mouhamadou SOW\nB3 Communication Digitale et Internationale\nFévrier à Avril 2026 - 3 mois"
    info_frame.text = info_text
    info_frame.paragraphs[0].font.size = Pt(18)
    info_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ==========================================
    # SLIDE 2 : PLAN DE LA PRÉSENTATION
    # ==========================================
    add_content_slide(
        "Plan de la Présentation",
        [
            "1. Introduction : Mon stage chez Buroxia",
            "2. Présentation de l'entreprise Buroxia",
            "3. Mes missions de Community Manager",
            "4. Mission phare : Collaboration avec Ayaa.lbns",
            "5. Bilan et perspectives professionnelles"
        ]
    )
    
    # ==========================================
    # SLIDE 3 : INTRODUCTION
    # ==========================================
    add_content_slide(
        "Introduction",
        [
            "🎓 Stage de 3 mois : Février à Avril 2026",
            "🏢 Entreprise : Buroxia (Prado Affaires), centre d'affaires à Marseille",
            "💼 Poste : Community Manager",
            "🎯 Objectif : Développer la visibilité digitale et la notoriété",
            "⭐ Mission phare : Partenariat avec l'influenceuse Ayaa.lbns"
        ]
    )
    
    # ==========================================
    # SLIDE 4 : BUROXIA - PRÉSENTATION
    # ==========================================
    add_content_slide(
        "Buroxia : Présentation de l'Entreprise",
        [
            "📍 Centre d'affaires situé 24 avenue du Prado, Marseille 6e",
            "📐 330 m² d'espaces professionnels",
            "📅 Plus de 15 ans d'expérience",
            "✅ Services : Domiciliation, bureaux équipés, salles de réunion",
            "👥 Clients : Entrepreneurs, TPE, professions libérales, influenceurs",
            "🎯 Positionnement : Adresse prestigieuse et relation de proximité"
        ]
    )
    
    # ==========================================
    # SLIDE 5 : BUROXIA - MARCHÉ
    # ==========================================
    add_content_slide(
        "Buroxia : Marché et Enjeux",
        [
            "🏙️ Marché marseillais des centres d'affaires très concurrentiel",
            "💡 Enjeu principal : Développer la visibilité et la notoriété",
            "📉 Problématique : Structure peu connue malgré une offre pertinente",
            "🚀 Solution : Stratégie digitale et marketing d'influence",
            "👤 Tuteur : Stephanne Heyoppe, propriétaire et directeur"
        ]
    )
    
    # ==========================================
    # SLIDE 6 : MES MISSIONS
    # ==========================================
    add_content_slide(
        "Mes Missions de Community Manager",
        [
            "📱 Gestion des réseaux sociaux (Instagram, TikTok, LinkedIn, Facebook)",
            "🎨 Création de contenus visuels et vidéo (Canva, montage)",
            "📄 Refonte de supports commerciaux (plaquette, grille tarifaire)",
            "🌐 Mise à jour du site web WordPress",
            "📞 Prospection commerciale et rédaction de devis",
            "🤝 Accompagnement client et organisation de réunions",
            "⭐ Mission phare : Marketing d'influence avec Ayaa.lbns"
        ]
    )
    
    # ==========================================
    # SLIDE 7 : TRANSITION MISSION PHARE
    # ==========================================
    add_highlight_slide(
        "Mission Phare",
        "Collaboration avec l'Influenceuse Ayaa.lbns",
        [
            "✨ Projet stratégique mené de A à Z",
            "🎯 De l'idée initiale à la signature du contrat",
            "💼 Responsabilité : Prospection, négociation, coordination"
        ]
    )
    
    # ==========================================
    # SLIDE 8 : GENÈSE DU PROJET
    # ==========================================
    add_content_slide(
        "Genèse du Projet d'Influence",
        [
            "💡 Initiative personnelle : Proposition lors d'une réunion avec la direction",
            "❓ Problématique : Comment accroître rapidement la visibilité de Buroxia ?",
            "🎯 Stratégie : S'appuyer sur le marketing d'influence",
            "✅ Validation : La direction valide le projet et me confie la mission",
            "📋 Objectif : Identifier et conclure un partenariat avec une influenceuse",
            "🚀 Responsabilité : Porteur du projet du début à la fin"
        ]
    )
    
    # ==========================================
    # SLIDE 9 : PROSPECTION
    # ==========================================
    add_content_slide(
        "Prospection d'Influenceurs",
        [
            "🔍 Veille active sur Instagram et TikTok",
            "📊 Analyse : audience, ligne éditoriale, adéquation avec Buroxia",
            "📧 Démarche : Messages privés et e-mails de contact",
            "❌ Difficultés : Plusieurs refus, manque de notoriété de Buroxia",
            "💪 Persévérance : Poursuite de la veille malgré les obstacles",
            "✨ Résultat : Identification d'une opportunité avec Ayaa.lbns"
        ]
    )
    
    # ==========================================
    # SLIDE 10 : AYAA.LBNS - PROFIL
    # ==========================================
    add_content_slide(
        "Profil de l'Influenceuse Ayaa.lbns",
        [
            "👤 Nom : Ouweik Chrystelle (pseudonyme : Ayaa.lbns)",
            "❤️ Engagement social : Aide aux personnes sans abri (France et international)",
            "🏠 Actions : Rénovation d'habitats, distribution de nourriture",
            "🌙 Visibilité accrue durant le Ramadan : ruptures du jeûne collectives",
            "👥 Impact : Réunit familles, entrepreneurs et personnes en difficulté",
            "✨ Image de marque : Authentique, engagée, humaine",
            "🤝 Adéquation : Valeurs de proximité alignées avec Buroxia"
        ]
    )
    
    # ==========================================
    # SLIDE 11 : L'OPPORTUNITÉ
    # ==========================================
    add_highlight_slide(
        "L'Opportunité en Or",
        "Repérage d'une Story Instagram",
        [
            "📱 Ayaa.lbns publie une story : recherche d'un bureau à Marseille pour 1 mois",
            "💡 Correspondance parfaite avec l'offre de Buroxia !",
            "📧 Contact immédiat avec le manager via l'e-mail indiqué",
            "⏰ Timing idéal pour proposer une collaboration structurée"
        ]
    )
    
    # ==========================================
    # SLIDE 12 : PREMIER CONTACT
    # ==========================================
    add_content_slide(
        "Premier Contact et Prise de Parole",
        [
            "📧 E-mail professionnel rédigé au manager d'Ayaa.lbns",
            "🏢 Présentation de Buroxia et de ses services",
            "💼 Exposition claire des solutions (domiciliation, bureaux, coworking)",
            "🤝 Ouverture vers un partenariat incluant de la visibilité",
            "✅ Résultat : Contact favorable, échange approfondi engagé"
        ]
    )
    
    # ==========================================
    # SLIDE 13 : OUTILS COMMERCIAUX
    # ==========================================
    add_content_slide(
        "Élaboration des Outils Commerciaux",
        [
            "📄 Fiche de prestations de services complète",
            "💰 Grille tarifaire détaillée (formules, promotions, conditions)",
            "🎨 Design professionnel et cohérent (Canva)",
            "💼 Utilisation : Base de négociation et présentations",
            "♻️ Réutilisation : Supports pour la prospection commerciale générale",
            "✅ Impact : Renforcement du sérieux de la proposition"
        ]
    )
    
    # ==========================================
    # SLIDE 14 : ORGANISATION RÉUNIONS
    # ==========================================
    add_content_slide(
        "Organisation des Réunions",
        [
            "🗓️ 3 réunions entre Buroxia et l'équipe d'Ayaa.lbns",
            "1️⃣ Réunion 1 (visio) : Présentation de l'offre et premiers échanges",
            "2️⃣ Réunion 2 (visio) : Affinage de la proposition et conditions",
            "3️⃣ Réunion 3 (présentiel) : Visite des locaux à Marseille + finalisation",
            "👥 Interlocuteurs : Direction Buroxia, Ayaa.lbns, son manager (Paris)",
            "📋 Mon rôle : Préparation des supports, participation active, suivi"
        ]
    )
    
    # ==========================================
    # SLIDE 15 : MODÈLE DE COLLABORATION
    # ==========================================
    add_highlight_slide(
        "Modèle de Collaboration Innovant",
        "Proposition Gagnant-Gagnant",
        [
            "💡 Clause spécifique : 50% prix location + 50% contrepartie publicitaire",
            "🏢 Pour Buroxia : Visibilité auprès d'une audience qualifiée",
            "👤 Pour Ayaa.lbns : Espace professionnel à conditions avantageuses",
            "✅ Résultat : Signature du contrat dans les locaux de Buroxia à Marseille"
        ]
    )
    
    # ==========================================
    # SLIDE 16 : RÉSULTATS
    # ==========================================
    add_content_slide(
        "Résultats et Indicateurs",
        [
            "✅ Signature d'une collaboration avec Ayaa.lbns",
            "📈 Exposition de Buroxia auprès de l'audience de l'influenceuse",
            "🏆 Association à une créatrice reconnue pour son engagement social",
            "📄 Livrables : Fiche de prestations, grille tarifaire, contrat signé",
            "🏢 Occupation : Mise à disposition d'un bureau pour 1 mois",
            "🚀 Impact : Premier partenariat influence structuré pour Buroxia"
        ]
    )
    
    # ==========================================
    # SLIDE 17 : VALEUR AJOUTÉE
    # ==========================================
    add_content_slide(
        "Valeur Ajoutée pour l'Entreprise",
        [
            "💼 Premier partenariat influenceur réussi pour Buroxia",
            "📊 Crédibilité renforcée auprès de prospects et partenaires",
            "👁️ Visibilité élargie sur une cible qualifiée",
            "🔄 Outils commerciaux réutilisables (fiche, grille tarifaire)",
            "🚪 Ouverture vers de futures collaborations d'influence",
            "🎯 Stratégie digitale structurée et opérationnelle"
        ]
    )
    
    # ==========================================
    # SLIDE 18 : ANALYSE CRITIQUE
    # ==========================================
    add_content_slide(
        "Analyse Critique et Axes d'Amélioration",
        [
            "✅ Points forts : Initiative, persévérance, opportunité saisie",
            "⚠️ Difficultés : Refus initiaux, faible notoriété, coordination à distance",
            "🔄 Amélioration possible : Processus de prospection systématisé",
            "📊 Recommandation : Mesurer les retombées (portée, engagement, leads)",
            "💡 Enseignement : Alignement des intérêts > Notoriété du profil"
        ]
    )
    
    # ==========================================
    # SLIDE 19 : COMPÉTENCES DÉVELOPPÉES
    # ==========================================
    add_content_slide(
        "Compétences Développées",
        [
            "🔍 Prospection de partenaires influenceurs (veille, ciblage, contact)",
            "📝 Rédaction de supports commerciaux structurés",
            "🎨 Création de contenus visuels et vidéo (Canva, montage)",
            "📅 Organisation et animation de réunions (visio et présentiel)",
            "💬 Négociation et construction d'offres gagnant-gagnant",
            "📈 Contribution à une stratégie de visibilité par l'influence",
            "🤝 Gestion de projet et coordination multi-parties"
        ]
    )
    
    # ==========================================
    # SLIDE 20 : SAVOIR-ÊTRE
    # ==========================================
    add_content_slide(
        "Savoir-Être et Posture Professionnelle",
        [
            "💪 Persévérance face aux refus",
            "💡 Initiative et prise de responsabilité",
            "✍️ Rigueur rédactionnelle et professionnelle",
            "🤝 Diplomatie en situation de négociation",
            "🎯 Capacité à représenter une entreprise avec crédibilité",
            "⚡ Réactivité et saisie d'opportunités en temps réel"
        ]
    )
    
    # ==========================================
    # SLIDE 21 : BILAN PERSONNEL
    # ==========================================
    add_content_slide(
        "Bilan Personnel",
        [
            "✅ Expérience formatrice et concrète",
            "🚀 Confirmation de mon intérêt pour la communication digitale",
            "🎯 Passion confirmée pour le marketing d'influence",
            "💼 Développement de compétences opérationnelles et stratégiques",
            "🤝 Travail en autonomie et collaboration avec la direction",
            "📈 Capacité à mener un projet stratégique de bout en bout"
        ]
    )
    
    # ==========================================
    # SLIDE 22 : PROJET PROFESSIONNEL
    # ==========================================
    add_highlight_slide(
        "Projet Professionnel",
        "Communication Digitale et Marketing d'Influence",
        [
            "🎯 Orientation : Stratégies de visibilité pour les marques",
            "💡 Objectif : Concevoir et piloter des projets d'influence",
            "🔗 Combinaison : Contenu créatif + Performance commerciale",
            "🚀 Ambition : Accompagner les marques dans leur transformation digitale"
        ]
    )
    
    # ==========================================
    # SLIDE 23 : CONCLUSION
    # ==========================================
    add_content_slide(
        "Conclusion",
        [
            "🏆 Mission phare réussie : De l'idée à la signature du contrat",
            "💼 Responsabilité assumée sur un projet stratégique",
            "📈 Contribution concrète au développement de Buroxia",
            "🎓 Compétences développées : prospection, négociation, création",
            "✨ Expérience valorisante pour mon projet professionnel",
            "🙏 Remerciements à Stephanne Heyoppe et à l'équipe Buroxia"
        ]
    )
    
    # ==========================================
    # SLIDE 24 : QUESTIONS
    # ==========================================
    slide_questions = add_title_slide(
        "Merci de votre attention",
        "Place aux questions"
    )
    
    # Sauvegarde
    prs.save('/workspace/Soutenance_Mouhamadou_Sow_B3_CDI.pptx')
    print("✅ Présentation créée avec succès : Soutenance_Mouhamadou_Sow_B3_CDI.pptx")
    print(f"📊 Nombre de slides : {len(prs.slides)}")
    print("🎯 Structure : Introduction → Buroxia → Missions → Ayaa.lbns → Bilan")

if __name__ == "__main__":
    create_presentation()
