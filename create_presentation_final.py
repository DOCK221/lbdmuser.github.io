#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Couleurs professionnelles
    COLOR_PRIMARY = RGBColor(0, 51, 102)  # Bleu marine professionnel
    COLOR_ACCENT = RGBColor(0, 112, 192)  # Bleu clair
    COLOR_TEXT = RGBColor(51, 51, 51)  # Gris foncé
    
    def add_title_slide(title, subtitle=""):
        """Slide de titre"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Fond bleu
        background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = COLOR_PRIMARY
        background.line.color.rgb = COLOR_PRIMARY
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Sous-titre
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
            tf2 = subtitle_box.text_frame
            tf2.text = subtitle
            p2 = tf2.paragraphs[0]
            p2.font.size = Pt(28)
            p2.font.color.rgb = RGBColor(255, 255, 255)
            p2.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_list):
        """Slide de contenu standard"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Fond blanc
        background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        background.line.color.rgb = RGBColor(255, 255, 255)
        
        # Barre de titre
        title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = COLOR_PRIMARY
        title_bar.line.color.rgb = COLOR_PRIMARY
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Contenu
        content_top = Inches(1.5)
        content_box = slide.shapes.add_textbox(Inches(1), content_top, Inches(8), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, text in enumerate(content_list):
            if i > 0:
                tf.add_paragraph()
            p = tf.paragraphs[i]
            p.text = text
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(16)
            p.level = 0
        
        return slide
    
    # ==========================================
    # SLIDE 1 : PAGE DE GARDE
    # ==========================================
    slide1 = add_title_slide(
        "Retour d'Expérience Professionnelle",
        "Stage Community Manager"
    )
    
    # Informations étudiant
    info_box = slide1.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1.5))
    tf = info_box.text_frame
    tf.text = "Mouhamadou SOW\nB3 Communication Digitale et Internationale\nBuroxia - Février à Avril 2026"
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # ==========================================
    # SLIDE 2 : PLAN
    # ==========================================
    add_content_slide(
        "Plan de la Présentation",
        [
            "1. Introduction : contexte et objectifs du stage",
            "",
            "2. Présentation de l'entreprise Buroxia",
            "",
            "3. Mes missions de Community Manager",
            "",
            "4. Mission phare : collaboration avec l'influenceuse Ayaa.lbns",
            "",
            "5. Bilan personnel et perspectives professionnelles"
        ]
    )
    
    # ==========================================
    # SLIDE 3 : INTRODUCTION
    # ==========================================
    add_content_slide(
        "Introduction",
        [
            "Stage de 3 mois effectué de février à avril 2026",
            "",
            "Entreprise : Buroxia, centre d'affaires situé à Marseille",
            "",
            "Poste : Community Manager",
            "",
            "Objectif principal : développer la visibilité digitale et la notoriété de Buroxia",
            "",
            "Mission phare : signature d'un partenariat avec l'influenceuse Ayaa.lbns"
        ]
    )
    
    # ==========================================
    # SLIDE 4 : BUROXIA PRÉSENTATION
    # ==========================================
    add_content_slide(
        "Présentation de Buroxia",
        [
            "Centre d'affaires implanté au 24 avenue du Prado, Marseille 6ème arrondissement",
            "",
            "330 m² d'espaces professionnels avec plus de 15 ans d'expérience",
            "",
            "Services proposés : domiciliation d'entreprises, location de bureaux équipés, mise à disposition de salles de réunion",
            "",
            "Clientèle cible : entrepreneurs, TPE, professions libérales et influenceurs",
            "",
            "Positionnement : adresse prestigieuse et relation de proximité avec les clients"
        ]
    )
    
    # ==========================================
    # SLIDE 5 : ENJEUX
    # ==========================================
    add_content_slide(
        "Enjeux et Contexte du Stage",
        [
            "Marché marseillais des centres d'affaires très concurrentiel",
            "",
            "Problématique principale : structure peu connue malgré une offre de qualité",
            "",
            "Enjeu stratégique : développer rapidement la visibilité et la notoriété",
            "",
            "Solution proposée : déployer une stratégie digitale couplée au marketing d'influence",
            "",
            "Tuteur de stage : Stephanne Heyoppe, propriétaire et directeur de Buroxia"
        ]
    )
    
    # ==========================================
    # SLIDE 6 : MES MISSIONS
    # ==========================================
    add_content_slide(
        "Mes Missions de Community Manager",
        [
            "Gestion quotidienne des réseaux sociaux : Instagram, TikTok, LinkedIn et Facebook",
            "",
            "Création de contenus visuels et vidéo avec Canva et outils de montage",
            "",
            "Refonte complète des supports commerciaux : plaquette, grille tarifaire et fiches de prestations",
            "",
            "Mise à jour et optimisation du site web WordPress de l'entreprise",
            "",
            "Prospection commerciale, rédaction de devis et accompagnement des clients"
        ]
    )
    
    # ==========================================
    # SLIDE 7 : MISSION PHARE
    # ==========================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond blanc
    bg7 = slide7.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg7.fill.solid()
    bg7.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg7.line.color.rgb = RGBColor(255, 255, 255)
    
    # Titre
    title7 = slide7.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.8))
    tf7 = title7.text_frame
    tf7.text = "Mission Phare"
    p7 = tf7.paragraphs[0]
    p7.font.size = Pt(40)
    p7.font.bold = True
    p7.font.color.rgb = COLOR_PRIMARY
    p7.alignment = PP_ALIGN.CENTER
    
    # Encadré
    box7 = slide7.shapes.add_shape(1, Inches(2), Inches(2.5), Inches(6), Inches(2))
    box7.fill.solid()
    box7.fill.fore_color.rgb = COLOR_ACCENT
    box7.line.color.rgb = COLOR_ACCENT
    
    # Texte encadré
    text7 = slide7.shapes.add_textbox(Inches(2.3), Inches(2.8), Inches(5.4), Inches(1.4))
    tf7b = text7.text_frame
    tf7b.text = "Collaboration avec l'influenceuse Ayaa.lbns"
    p7b = tf7b.paragraphs[0]
    p7b.font.size = Pt(28)
    p7b.font.bold = True
    p7b.font.color.rgb = RGBColor(255, 255, 255)
    p7b.alignment = PP_ALIGN.CENTER
    
    # Détails
    details7 = slide7.shapes.add_textbox(Inches(1.5), Inches(5), Inches(7), Inches(1.5))
    tf7c = details7.text_frame
    details_text = "Projet stratégique mené de A à Z : de l'idée initiale à la signature du contrat"
    tf7c.text = details_text
    p7c = tf7c.paragraphs[0]
    p7c.font.size = Pt(20)
    p7c.font.color.rgb = COLOR_TEXT
    p7c.alignment = PP_ALIGN.CENTER
    
    # ==========================================
    # SLIDE 8 : GENÈSE
    # ==========================================
    add_content_slide(
        "Genèse du Projet d'Influence",
        [
            "Initiative personnelle présentée lors d'une réunion avec la direction",
            "",
            "Problématique identifiée : comment accroître rapidement la visibilité de Buroxia dans un marché concurrentiel ?",
            "",
            "Stratégie proposée : utiliser le marketing d'influence pour toucher une audience qualifiée",
            "",
            "Validation du projet par la direction qui me confie la mission complète",
            "",
            "Objectif : identifier, contacter et conclure un partenariat avec une influenceuse pertinente"
        ]
    )
    
    # ==========================================
    # SLIDE 9 : PROSPECTION
    # ==========================================
    add_content_slide(
        "Prospection d'Influenceurs",
        [
            "Mise en place d'une veille active quotidienne sur Instagram et TikTok",
            "",
            "Analyse approfondie de chaque profil : audience, ligne éditoriale, valeurs et adéquation avec Buroxia",
            "",
            "Prise de contact par messages privés et e-mails professionnels",
            "",
            "Difficultés rencontrées : plusieurs refus liés au manque de notoriété de Buroxia",
            "",
            "Persévérance maintenue jusqu'à l'identification de l'opportunité avec Ayaa.lbns"
        ]
    )
    
    # ==========================================
    # SLIDE 10 : PROFIL AYAA
    # ==========================================
    add_content_slide(
        "Profil de l'Influenceuse Ayaa.lbns",
        [
            "Nom : Ouweik Chrystelle, connue sous le pseudonyme Ayaa.lbns",
            "",
            "Engagement social fort : aide régulière aux personnes sans abri en France et à l'international",
            "",
            "Actions concrètes : rénovation d'habitats, distribution de nourriture financée personnellement",
            "",
            "Visibilité accrue durant le Ramadan avec organisation de ruptures du jeûne collectives",
            "",
            "Image de marque authentique et humaine, parfaitement alignée avec les valeurs de proximité de Buroxia"
        ]
    )
    
    # ==========================================
    # SLIDE 11 : OPPORTUNITÉ
    # ==========================================
    add_content_slide(
        "L'Opportunité Saisie",
        [
            "Repérage d'une story Instagram où Ayaa.lbns recherche un bureau à Marseille pour un mois",
            "",
            "Correspondance parfaite avec l'offre de Buroxia : espaces professionnels flexibles, localisation centrale",
            "",
            "Contact immédiat du manager via l'adresse e-mail indiquée dans la story",
            "",
            "Timing idéal pour proposer une collaboration structurée plutôt qu'un simple contact commercial"
        ]
    )
    
    # ==========================================
    # SLIDE 12 : PREMIER CONTACT
    # ==========================================
    add_content_slide(
        "Premier Contact Professionnel",
        [
            "Rédaction d'un e-mail professionnel adressé au manager d'Ayaa.lbns",
            "",
            "Présentation détaillée de Buroxia et de l'ensemble de ses services",
            "",
            "Mise en avant des solutions adaptées : domiciliation, bureaux équipés, espaces de coworking",
            "",
            "Ouverture vers un partenariat incluant une dimension de visibilité pour Buroxia",
            "",
            "Résultat : réponse favorable et ouverture d'échanges approfondis"
        ]
    )
    
    # ==========================================
    # SLIDE 13 : OUTILS COMMERCIAUX
    # ==========================================
    add_content_slide(
        "Élaboration des Outils Commerciaux",
        [
            "Création d'une fiche de prestations de services complète et détaillée",
            "",
            "Conception d'une grille tarifaire claire : formules, promotions et conditions",
            "",
            "Design professionnel et cohérent réalisé avec Canva",
            "",
            "Ces documents ont servi de base à la négociation et aux présentations lors des réunions",
            "",
            "Réutilisation de ces supports pour l'ensemble de la prospection commerciale de Buroxia"
        ]
    )
    
    # ==========================================
    # SLIDE 14 : RÉUNIONS
    # ==========================================
    add_content_slide(
        "Organisation des Réunions",
        [
            "Trois réunions organisées entre l'équipe de Buroxia et l'équipe d'Ayaa.lbns",
            "",
            "Première réunion en visioconférence : présentation de l'offre et premiers échanges",
            "",
            "Deuxième réunion en visioconférence : affinage de la proposition et discussion des conditions",
            "",
            "Troisième réunion en présentiel à Marseille : visite des locaux et finalisation du partenariat",
            "",
            "Mon rôle : préparation des supports, participation active aux échanges et suivi opérationnel complet"
        ]
    )
    
    # ==========================================
    # SLIDE 15 : MODÈLE
    # ==========================================
    add_content_slide(
        "Modèle de Collaboration Innovant",
        [
            "Proposition d'une clause spécifique : 50% du prix correspond à la location, 50% à la contrepartie publicitaire",
            "",
            "Avantages pour Buroxia : visibilité auprès d'une audience qualifiée et accueil d'une créatrice de contenu",
            "",
            "Avantages pour Ayaa.lbns : espace professionnel à conditions avantageuses et valorisation du lieu",
            "",
            "Modèle gagnant-gagnant qui a permis la signature du contrat dans les locaux de Buroxia"
        ]
    )
    
    # ==========================================
    # SLIDE 16 : RÉSULTATS
    # ==========================================
    add_content_slide(
        "Résultats Obtenus",
        [
            "Signature effective d'une collaboration avec l'influenceuse Ayaa.lbns",
            "",
            "Exposition de la marque Buroxia auprès de l'audience engagée de l'influenceuse",
            "",
            "Association de Buroxia à une créatrice reconnue pour son engagement social",
            "",
            "Livrables produits : fiche de prestations, grille tarifaire, supports de réunion et contrat signé",
            "",
            "Premier partenariat d'influence structuré pour Buroxia, ouvrant la voie à de futures collaborations"
        ]
    )
    
    # ==========================================
    # SLIDE 17 : VALEUR AJOUTÉE
    # ==========================================
    add_content_slide(
        "Valeur Ajoutée pour l'Entreprise",
        [
            "Concrétisation du premier partenariat influenceur pour Buroxia",
            "",
            "Renforcement de la crédibilité de la marque auprès des prospects et partenaires",
            "",
            "Élargissement de la visibilité sur une cible qualifiée",
            "",
            "Création d'outils commerciaux réutilisables pour l'ensemble de la prospection",
            "",
            "Ouverture stratégique vers de futures collaborations d'influence"
        ]
    )
    
    # ==========================================
    # SLIDE 18 : ANALYSE CRITIQUE
    # ==========================================
    add_content_slide(
        "Analyse Critique et Axes d'Amélioration",
        [
            "Points forts : initiative personnelle, persévérance face aux refus, saisie de l'opportunité",
            "",
            "Difficultés rencontrées : refus initiaux, faible notoriété de Buroxia, coordination à distance avec l'équipe basée à Paris",
            "",
            "Axes d'amélioration : systématiser le processus de prospection avec création d'une base de données d'influenceurs",
            "",
            "Recommandation : mettre en place des indicateurs de mesure des retombées (portée, engagement, leads générés)"
        ]
    )
    
    # ==========================================
    # SLIDE 19 : COMPÉTENCES
    # ==========================================
    add_content_slide(
        "Compétences Développées",
        [
            "Prospection de partenaires influenceurs par veille active, ciblage et prises de contact",
            "",
            "Rédaction de supports commerciaux structurés et professionnels",
            "",
            "Création de contenus visuels et vidéo pour les réseaux sociaux et supports print",
            "",
            "Organisation et animation de réunions à distance et en présentiel",
            "",
            "Négociation commerciale et construction d'offres gagnant-gagnant",
            "",
            "Gestion complète de projet avec coordination de multiples parties prenantes"
        ]
    )
    
    # ==========================================
    # SLIDE 20 : SAVOIR-ÊTRE
    # ==========================================
    add_content_slide(
        "Savoir-Être Développé",
        [
            "Persévérance et résilience face aux refus et aux obstacles",
            "",
            "Initiative et capacité à proposer des projets stratégiques",
            "",
            "Prise de responsabilité sur un projet de bout en bout",
            "",
            "Rigueur rédactionnelle et professionnalisme dans les échanges",
            "",
            "Diplomatie et capacité de négociation en contexte professionnel",
            "",
            "Réactivité et saisie d'opportunités en temps réel"
        ]
    )
    
    # ==========================================
    # SLIDE 21 : BILAN PERSONNEL
    # ==========================================
    add_content_slide(
        "Bilan Personnel",
        [
            "Expérience professionnelle extrêmement formatrice et concrète",
            "",
            "Confirmation de mon intérêt pour la communication digitale et le marketing d'influence",
            "",
            "Développement de compétences à la fois opérationnelles et stratégiques",
            "",
            "Capacité démontrée à travailler en autonomie tout en collaborant avec la direction",
            "",
            "Réussite d'un projet stratégique mené de l'idée initiale jusqu'à la signature du contrat"
        ]
    )
    
    # ==========================================
    # SLIDE 22 : PROJET PRO
    # ==========================================
    add_content_slide(
        "Projet Professionnel",
        [
            "Orientation confirmée vers la communication digitale et le marketing d'influence",
            "",
            "Objectif : concevoir et piloter des stratégies de visibilité pour les marques",
            "",
            "Ambition : combiner créativité dans la création de contenu et performance commerciale",
            "",
            "Souhait d'accompagner les marques dans leur transformation digitale et leur stratégie d'influence"
        ]
    )
    
    # ==========================================
    # SLIDE 23 : CONCLUSION
    # ==========================================
    add_content_slide(
        "Conclusion",
        [
            "Mission phare réussie : de l'idée initiale à la signature du contrat avec Ayaa.lbns",
            "",
            "Responsabilité assumée sur un projet stratégique rare pour un stagiaire",
            "",
            "Contribution concrète et mesurable au développement de Buroxia",
            "",
            "Compétences développées en prospection, négociation, création et gestion de projet",
            "",
            "Expérience valorisante qui confirme mon projet professionnel",
            "",
            "Remerciements à Stephanne Heyoppe et à toute l'équipe Buroxia pour leur confiance"
        ]
    )
    
    # ==========================================
    # SLIDE 24 : QUESTIONS
    # ==========================================
    slide24 = add_title_slide(
        "Merci de votre attention",
        "Je suis à votre disposition pour vos questions"
    )
    
    # Sauvegarde
    prs.save('/workspace/Soutenance_Mouhamadou_Sow_Final.pptx')
    print("✅ Présentation finale créée : Soutenance_Mouhamadou_Sow_Final.pptx")
    print(f"📊 Nombre de slides : {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
