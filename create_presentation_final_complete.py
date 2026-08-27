#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Palette de couleurs
    COLOR_PRIMARY = RGBColor(0, 102, 204)
    COLOR_SECONDARY = RGBColor(255, 153, 0)
    COLOR_ACCENT = RGBColor(51, 51, 51)
    COLOR_LIGHT = RGBColor(240, 248, 255)
    
    # Chemins des images
    IMG_OFFRES = "/home/ubuntu/.cursor/projects/workspace/assets/01a044c7-8d47-75c0-877b-3348737546ad.jpg"
    IMG_BUREAU1 = "/home/ubuntu/.cursor/projects/workspace/assets/01a044c7-8e03-76fa-a38e-821e29068e98.jpg"
    IMG_BUREAU2 = "/home/ubuntu/.cursor/projects/workspace/assets/01a044c7-8e6f-7b1b-97ea-f2c918a20bbf.jpg"
    IMG_SERVICES = "/home/ubuntu/.cursor/projects/workspace/assets/01a044c7-8e86-7cb4-9bef-6cfd84137737.jpg"
    IMG_BUREAU406 = "/home/ubuntu/.cursor/projects/workspace/assets/01a044c7-8e9e-77c8-82cb-5d3abf3ef9c2.jpg"
    
    def add_title_slide(main_title, subtitle, student_info):
        """Page de garde"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_PRIMARY
        bg.line.color.rgb = COLOR_PRIMARY
        
        circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0), Inches(1.5), Inches(1.5))
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = RGBColor(0, 82, 184)
        circle1.line.color.rgb = RGBColor(0, 82, 184)
        
        circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0), Inches(6), Inches(1.5), Inches(1.5))
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = RGBColor(0, 82, 184)
        circle2.line.color.rgb = RGBColor(0, 82, 184)
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.8))
        tf = title_box.text_frame
        tf.text = main_title
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.1
        
        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.6))
        tf2 = subtitle_box.text_frame
        tf2.text = subtitle
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
        
        info_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.4), Inches(6), Inches(1.4))
        info_bg.fill.solid()
        info_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        info_bg.line.color.rgb = RGBColor(255, 255, 255)
        
        info_box = slide.shapes.add_textbox(Inches(2.2), Inches(5.5), Inches(5.6), Inches(1.2))
        tf3 = info_box.text_frame
        tf3.text = student_info
        tf3.word_wrap = True
        tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in tf3.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = COLOR_ACCENT
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.space_after = Pt(4)
        
        return slide
    
    def add_content_slide(title, content_list, slide_number=None):
        """Slide de contenu standard"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(255, 255, 255)
        
        side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = COLOR_PRIMARY
        side_bar.line.color.rgb = COLOR_PRIMARY
        
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), 0, Inches(9.75), Inches(1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_LIGHT
        header_shape.line.color.rgb = COLOR_LIGHT
        
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.95), Inches(9.75), Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLOR_SECONDARY
        accent_line.line.color.rgb = COLOR_SECONDARY
        
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.25), Inches(8.5), Inches(0.65))
        tf = title_box.text_frame
        tf.text = title
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        if slide_number:
            num_box = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.7), Inches(0.3))
            tf_num = num_box.text_frame
            tf_num.text = str(slide_number)
            p_num = tf_num.paragraphs[0]
            p_num.font.size = Pt(12)
            p_num.font.color.rgb = COLOR_PRIMARY
            p_num.alignment = PP_ALIGN.RIGHT
        
        content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = COLOR_LIGHT
        content_bg.line.color.rgb = COLOR_LIGHT
        
        content_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(7.6), Inches(4.9))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        
        for i, text in enumerate(content_list):
            if i > 0:
                tf.add_paragraph()
            p = tf.paragraphs[i]
            
            if text.strip():
                p.text = "• " + text
                p.font.size = Pt(17)
                p.font.color.rgb = COLOR_ACCENT
                p.space_after = Pt(10)
                p.line_spacing = 1.15
            else:
                p.text = ""
                p.space_after = Pt(4)
        
        return slide
    
    def add_content_slide_with_image(title, content_list, image_path, slide_number=None):
        """Slide de contenu avec image"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(255, 255, 255)
        
        side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = COLOR_PRIMARY
        side_bar.line.color.rgb = COLOR_PRIMARY
        
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), 0, Inches(9.75), Inches(1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_LIGHT
        header_shape.line.color.rgb = COLOR_LIGHT
        
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.95), Inches(9.75), Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLOR_SECONDARY
        accent_line.line.color.rgb = COLOR_SECONDARY
        
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.25), Inches(8.5), Inches(0.65))
        tf = title_box.text_frame
        tf.text = title
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        if slide_number:
            num_box = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.7), Inches(0.3))
            tf_num = num_box.text_frame
            tf_num.text = str(slide_number)
            p_num = tf_num.paragraphs[0]
            p_num.font.size = Pt(12)
            p_num.font.color.rgb = COLOR_PRIMARY
            p_num.alignment = PP_ALIGN.RIGHT
        
        # Zone de contenu texte à gauche
        content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(4.5), Inches(5.5))
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = COLOR_LIGHT
        content_bg.line.color.rgb = COLOR_LIGHT
        
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(4), Inches(4.9))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        
        for i, text in enumerate(content_list):
            if i > 0:
                tf.add_paragraph()
            p = tf.paragraphs[i]
            
            if text.strip():
                p.text = "• " + text
                p.font.size = Pt(15)
                p.font.color.rgb = COLOR_ACCENT
                p.space_after = Pt(8)
                p.line_spacing = 1.15
            else:
                p.text = ""
                p.space_after = Pt(3)
        
        # Image à droite
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.8), width=Inches(3.8))
        
        return slide
    
    def add_visual_examples_slide(slide_number):
        """Slide spéciale pour les exemples de créations visuelles"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(255, 255, 255)
        
        side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = COLOR_PRIMARY
        side_bar.line.color.rgb = COLOR_PRIMARY
        
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), 0, Inches(9.75), Inches(1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_LIGHT
        header_shape.line.color.rgb = COLOR_LIGHT
        
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0.95), Inches(9.75), Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLOR_SECONDARY
        accent_line.line.color.rgb = COLOR_SECONDARY
        
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.25), Inches(8.5), Inches(0.65))
        tf = title_box.text_frame
        tf.text = "Exemples de Créations Visuelles"
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        num_box = slide.shapes.add_textbox(Inches(9), Inches(7.1), Inches(0.7), Inches(0.3))
        tf_num = num_box.text_frame
        tf_num.text = str(slide_number)
        p_num = tf_num.paragraphs[0]
        p_num.font.size = Pt(12)
        p_num.font.color.rgb = COLOR_PRIMARY
        p_num.alignment = PP_ALIGN.RIGHT
        
        # Image 1 - Offres de domiciliation (en haut gauche)
        if os.path.exists(IMG_OFFRES):
            slide.shapes.add_picture(IMG_OFFRES, Inches(0.75), Inches(1.5), width=Inches(4.4))
        
        # Titre image 1
        title1_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.3), Inches(4.4), Inches(0.4))
        tf1 = title1_box.text_frame
        tf1.text = "Plaquette Offres de Domiciliation"
        tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf1.paragraphs[0]
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_PRIMARY
        p1.alignment = PP_ALIGN.CENTER
        
        # Image 2 - Services d'entreprise (en haut droite)
        if os.path.exists(IMG_SERVICES):
            slide.shapes.add_picture(IMG_SERVICES, Inches(5.35), Inches(1.5), width=Inches(3.9))
        
        # Titre image 2
        title2_box = slide.shapes.add_textbox(Inches(5.35), Inches(4.3), Inches(3.9), Inches(0.4))
        tf2 = title2_box.text_frame
        tf2.text = "Infographie Services Buroxia"
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_PRIMARY
        p2.alignment = PP_ALIGN.CENTER
        
        # Image 3 - Bureau 406 (en bas)
        if os.path.exists(IMG_BUREAU406):
            slide.shapes.add_picture(IMG_BUREAU406, Inches(3), Inches(4.9), width=Inches(4))
        
        # Titre image 3
        title3_box = slide.shapes.add_textbox(Inches(3), Inches(6.8), Inches(4), Inches(0.4))
        tf3 = title3_box.text_frame
        tf3.text = "Publication Bureau 406 - Coworking"
        tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
        p3 = tf3.paragraphs[0]
        p3.font.size = Pt(14)
        p3.font.bold = True
        p3.font.color.rgb = COLOR_PRIMARY
        p3.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_highlight_slide(title, highlight_text, details_list):
        """Slide de mise en avant"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(255, 255, 255)
        
        side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = COLOR_PRIMARY
        side_bar.line.color.rgb = COLOR_PRIMARY
        
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(7), Inches(0.8))
        tf = title_box.text_frame
        tf.text = title
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.65), Inches(2.35), Inches(6.7), Inches(1.8))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
        shadow.line.color.rgb = RGBColor(200, 200, 200)
        
        highlight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(6.7), Inches(1.8))
        highlight_box.fill.solid()
        highlight_box.fill.fore_color.rgb = COLOR_SECONDARY
        highlight_box.line.color.rgb = COLOR_SECONDARY
        
        text_box = slide.shapes.add_textbox(Inches(1.8), Inches(2.5), Inches(6.1), Inches(1.2))
        tf2 = text_box.text_frame
        tf2.text = highlight_text
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
        p2.line_spacing = 1.1
        
        details_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(2.2))
        tf3 = details_box.text_frame
        tf3.word_wrap = True
        tf3.vertical_anchor = MSO_ANCHOR.TOP
        
        for i, detail in enumerate(details_list):
            if i > 0:
                tf3.add_paragraph()
            p3 = tf3.paragraphs[i]
            p3.text = "• " + detail
            p3.font.size = Pt(17)
            p3.font.color.rgb = COLOR_ACCENT
            p3.space_after = Pt(8)
            p3.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # ==========================================
    # CRÉATION DES SLIDES
    # ==========================================
    
    # SLIDE 1
    add_title_slide(
        "Retour d'Expérience\nProfessionnelle",
        "Stage Community Manager",
        "Mouhamadou SOW\nB3 Communication Digitale et Internationale\nBuroxia - Février à Avril 2026"
    )
    
    # SLIDE 2
    add_content_slide(
        "Plan de la Présentation",
        [
            "Introduction : contexte et objectifs du stage",
            "",
            "Présentation de l'entreprise Buroxia",
            "",
            "Mes missions de Community Manager",
            "",
            "Mission phare : collaboration avec l'influenceuse Ayaa.lbns",
            "",
            "Bilan personnel et perspectives professionnelles"
        ],
        2
    )
    
    # SLIDE 3
    add_content_slide(
        "Introduction",
        [
            "Stage de 3 mois effectué de février à avril 2026",
            "",
            "Entreprise : Buroxia, centre d'affaires situé à Marseille",
            "",
            "Poste : Community Manager",
            "",
            "Objectif principal : développer la visibilité digitale et la notoriété de Buroxia",
            "",
            "Mission phare : signature d'un partenariat avec l'influenceuse Ayaa.lbns"
        ],
        3
    )
    
    # SLIDE 4 - AVEC IMAGE DES BUREAUX
    add_content_slide_with_image(
        "Présentation de Buroxia",
        [
            "Centre d'affaires implanté au 24 avenue du Prado, Marseille 6ème",
            "",
            "330 m² d'espaces professionnels avec plus de 15 ans d'expérience",
            "",
            "Services : domiciliation, bureaux équipés, salles de réunion",
            "",
            "Clientèle : entrepreneurs, TPE, professions libérales, influenceurs"
        ],
        IMG_BUREAU1,
        4
    )
    
    # SLIDE 5
    add_content_slide(
        "Enjeux et Contexte du Stage",
        [
            "Marché marseillais des centres d'affaires très concurrentiel",
            "",
            "Problématique : structure peu connue malgré une offre de qualité",
            "",
            "Enjeu stratégique : développer rapidement la visibilité et la notoriété",
            "",
            "Solution : stratégie digitale couplée au marketing d'influence",
            "",
            "Tuteur : Stephanne Heyoppe, propriétaire et directeur de Buroxia"
        ],
        5
    )
    
    # SLIDE 6 - AVEC IMAGE BUREAU 2
    add_content_slide_with_image(
        "Mes Missions de Community Manager",
        [
            "Gestion des réseaux sociaux : Instagram, TikTok, LinkedIn, Facebook",
            "",
            "Création de contenus visuels et vidéo",
            "",
            "Refonte des supports commerciaux",
            "",
            "Mise à jour du site web WordPress",
            "",
            "Prospection commerciale et accompagnement clients"
        ],
        IMG_BUREAU2,
        6
    )
    
    # SLIDE 7
    add_highlight_slide(
        "Mission Phare",
        "Collaboration avec l'influenceuse Ayaa.lbns",
        [
            "Projet stratégique mené de A à Z",
            "De l'idée initiale à la signature du contrat",
            "Responsabilité : prospection, négociation, coordination"
        ]
    )
    
    # SLIDE 8
    add_content_slide(
        "Genèse du Projet d'Influence",
        [
            "Initiative personnelle présentée lors d'une réunion avec la direction",
            "",
            "Problématique identifiée : comment accroître rapidement la visibilité de Buroxia ?",
            "",
            "Stratégie proposée : utiliser le marketing d'influence pour toucher une audience qualifiée",
            "",
            "Validation du projet par la direction qui me confie la mission complète",
            "",
            "Objectif : identifier, contacter et conclure un partenariat avec une influenceuse"
        ],
        8
    )
    
    # SLIDE 9
    add_content_slide(
        "Prospection d'Influenceurs",
        [
            "Mise en place d'une veille active quotidienne sur Instagram et TikTok",
            "",
            "Analyse de chaque profil : audience, ligne éditoriale, valeurs, adéquation",
            "",
            "Prises de contact par messages privés et e-mails professionnels",
            "",
            "Difficultés : plusieurs refus liés au manque de notoriété de Buroxia",
            "",
            "Persévérance maintenue jusqu'à l'opportunité avec Ayaa.lbns"
        ],
        9
    )
    
    # SLIDE 10
    add_content_slide(
        "Profil de l'Influenceuse Ayaa.lbns",
        [
            "Nom : Ouweik Chrystelle, pseudonyme Ayaa.lbns",
            "",
            "Engagement social fort : aide aux personnes sans abri en France et à l'international",
            "",
            "Actions concrètes : rénovation d'habitats, distribution de nourriture",
            "",
            "Visibilité accrue durant le Ramadan avec ruptures du jeûne collectives",
            "",
            "Image authentique et humaine, alignée avec les valeurs de proximité de Buroxia"
        ],
        10
    )
    
    # SLIDE 11
    add_content_slide(
        "L'Opportunité Saisie",
        [
            "Repérage d'une story Instagram : Ayaa.lbns recherche un bureau à Marseille",
            "",
            "Correspondance parfaite avec l'offre de Buroxia",
            "",
            "Contact immédiat du manager via l'adresse e-mail indiquée",
            "",
            "Timing idéal pour proposer une collaboration structurée"
        ],
        11
    )
    
    # SLIDE 12
    add_content_slide(
        "Premier Contact Professionnel",
        [
            "Rédaction d'un e-mail professionnel au manager d'Ayaa.lbns",
            "",
            "Présentation détaillée de Buroxia et de ses services",
            "",
            "Mise en avant des solutions : domiciliation, bureaux, coworking",
            "",
            "Ouverture vers un partenariat incluant de la visibilité pour Buroxia",
            "",
            "Résultat : réponse favorable et ouverture d'échanges approfondis"
        ],
        12
    )
    
    # SLIDE 13
    add_content_slide(
        "Élaboration des Outils Commerciaux",
        [
            "Création d'une fiche de prestations de services complète",
            "",
            "Conception d'une grille tarifaire claire avec formules et promotions",
            "",
            "Design professionnel et cohérent réalisé avec Canva",
            "",
            "Utilisation pour la négociation et les présentations lors des réunions",
            "",
            "Réutilisation pour l'ensemble de la prospection commerciale de Buroxia"
        ],
        13
    )
    
    # SLIDE 14
    add_content_slide(
        "Organisation des Réunions",
        [
            "Trois réunions entre l'équipe Buroxia et l'équipe d'Ayaa.lbns",
            "",
            "Réunion 1 (visio) : présentation de l'offre et premiers échanges",
            "",
            "Réunion 2 (visio) : affinage de la proposition et discussions",
            "",
            "Réunion 3 (présentiel) : visite des locaux et finalisation du partenariat",
            "",
            "Mon rôle : préparation des supports, participation active et suivi complet"
        ],
        14
    )
    
    # SLIDE 15
    add_content_slide(
        "Modèle de Collaboration Innovant",
        [
            "Clause spécifique : 50 % location + 50 % contrepartie publicitaire",
            "",
            "Pour Buroxia : visibilité auprès d'une audience qualifiée",
            "",
            "Pour Ayaa.lbns : espace professionnel à conditions avantageuses",
            "",
            "Modèle gagnant-gagnant ayant permis la signature du contrat"
        ],
        15
    )
    
    # SLIDE 16
    add_content_slide(
        "Résultats Obtenus",
        [
            "Signature effective d'une collaboration avec Ayaa.lbns",
            "",
            "Exposition de Buroxia auprès de l'audience engagée de l'influenceuse",
            "",
            "Association à une créatrice reconnue pour son engagement social",
            "",
            "Livrables : fiche de prestations, grille tarifaire, supports, contrat",
            "",
            "Premier partenariat d'influence structuré ouvrant la voie à d'autres"
        ],
        16
    )
    
    # SLIDE 17
    add_content_slide(
        "Valeur Ajoutée pour l'Entreprise",
        [
            "Concrétisation du premier partenariat influenceur pour Buroxia",
            "",
            "Renforcement de la crédibilité de la marque",
            "",
            "Élargissement de la visibilité sur une cible qualifiée",
            "",
            "Création d'outils commerciaux réutilisables",
            "",
            "Ouverture stratégique vers de futures collaborations d'influence"
        ],
        17
    )
    
    # SLIDE 18
    add_content_slide(
        "Analyse Critique et Axes d'Amélioration",
        [
            "Points forts : initiative personnelle, persévérance, saisie de l'opportunité",
            "",
            "Difficultés : refus initiaux, faible notoriété, coordination à distance",
            "",
            "Axes d'amélioration : systématiser le processus de prospection",
            "",
            "Recommandation : mettre en place des indicateurs de mesure des retombées"
        ],
        18
    )
    
    # SLIDE 19
    add_content_slide(
        "Compétences Développées",
        [
            "Prospection de partenaires influenceurs par veille active et ciblage",
            "",
            "Rédaction de supports commerciaux structurés et professionnels",
            "",
            "Création de contenus visuels et vidéo pour réseaux sociaux",
            "",
            "Organisation et animation de réunions à distance et en présentiel",
            "",
            "Négociation commerciale et construction d'offres gagnant-gagnant"
        ],
        19
    )
    
    # SLIDE 20
    add_content_slide(
        "Savoir-Être Développé",
        [
            "Persévérance et résilience face aux refus et aux obstacles",
            "",
            "Initiative et capacité à proposer des projets stratégiques",
            "",
            "Prise de responsabilité sur un projet de bout en bout",
            "",
            "Rigueur rédactionnelle et professionnalisme dans les échanges",
            "",
            "Réactivité et saisie d'opportunités en temps réel"
        ],
        20
    )
    
    # SLIDE 21
    add_content_slide(
        "Bilan Personnel",
        [
            "Expérience professionnelle extrêmement formatrice et concrète",
            "",
            "Confirmation de mon intérêt pour la communication digitale",
            "",
            "Développement de compétences opérationnelles et stratégiques",
            "",
            "Capacité à travailler en autonomie avec collaboration directionnelle",
            "",
            "Réussite d'un projet stratégique de l'idée à la signature"
        ],
        21
    )
    
    # SLIDE 22
    add_content_slide(
        "Projet Professionnel",
        [
            "Orientation confirmée : communication digitale et marketing d'influence",
            "",
            "Objectif : concevoir et piloter des stratégies de visibilité pour les marques",
            "",
            "Ambition : combiner créativité et performance commerciale",
            "",
            "Souhait : accompagner les marques dans leur transformation digitale"
        ],
        22
    )
    
    # SLIDE 23 - NOUVELLE SLIDE : EXEMPLES DE CRÉATIONS VISUELLES (AVANT LA CONCLUSION)
    add_visual_examples_slide(23)
    
    # SLIDE 24 - CONCLUSION
    add_content_slide(
        "Conclusion",
        [
            "Mission phare réussie : de l'idée initiale à la signature du contrat",
            "",
            "Responsabilité assumée sur un projet stratégique rare pour un stagiaire",
            "",
            "Contribution concrète et mesurable au développement de Buroxia",
            "",
            "Compétences développées en prospection, négociation et création",
            "",
            "Remerciements à Stephanne Heyoppe et à toute l'équipe Buroxia"
        ],
        24
    )
    
    # SLIDE 25 - QUESTIONS
    slide25 = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg25 = slide25.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg25.fill.solid()
    bg25.fill.fore_color.rgb = COLOR_PRIMARY
    bg25.line.color.rgb = COLOR_PRIMARY
    
    circle25_1 = slide25.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0), Inches(1.5), Inches(1.5))
    circle25_1.fill.solid()
    circle25_1.fill.fore_color.rgb = RGBColor(0, 82, 184)
    circle25_1.line.color.rgb = RGBColor(0, 82, 184)
    
    circle25_2 = slide25.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0), Inches(6), Inches(1.5), Inches(1.5))
    circle25_2.fill.solid()
    circle25_2.fill.fore_color.rgb = RGBColor(0, 82, 184)
    circle25_2.line.color.rgb = RGBColor(0, 82, 184)
    
    main_text25 = slide25.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(0.9))
    tf25 = main_text25.text_frame
    tf25.text = "Merci de votre attention"
    tf25.vertical_anchor = MSO_ANCHOR.MIDDLE
    p25 = tf25.paragraphs[0]
    p25.font.size = Pt(48)
    p25.font.bold = True
    p25.font.color.rgb = RGBColor(255, 255, 255)
    p25.alignment = PP_ALIGN.CENTER
    
    sub_text25 = slide25.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(0.7))
    tf25b = sub_text25.text_frame
    tf25b.text = "Je suis à votre disposition pour vos questions"
    tf25b.vertical_anchor = MSO_ANCHOR.MIDDLE
    p25b = tf25b.paragraphs[0]
    p25b.font.size = Pt(26)
    p25b.font.color.rgb = RGBColor(255, 255, 255)
    p25b.alignment = PP_ALIGN.CENTER
    
    # Sauvegarde
    prs.save('/workspace/Soutenance_Mouhamadou_Sow_FINAL_COMPLETE.pptx')
    print("✅ Présentation FINALE COMPLÈTE créée !")
    print("📸 Toutes les 5 images intégrées")
    print("🎨 Nouvelle slide 23 : Exemples de Créations Visuelles")
    print(f"📊 {len(prs.slides)} slides au total")

if __name__ == "__main__":
    create_presentation()
