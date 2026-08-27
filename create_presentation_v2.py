#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Palette de couleurs moderne
    COLOR_PRIMARY = RGBColor(0, 102, 204)  # Bleu vif
    COLOR_SECONDARY = RGBColor(255, 153, 0)  # Orange
    COLOR_ACCENT = RGBColor(51, 51, 51)  # Gris foncé
    COLOR_LIGHT = RGBColor(240, 248, 255)  # Bleu très clair
    
    def add_title_slide(main_title, subtitle, student_info):
        """Slide de titre avec design moderne"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Fond dégradé simulé avec formes
        bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = COLOR_PRIMARY
        bg1.line.color.rgb = COLOR_PRIMARY
        
        # Formes décoratives
        circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(0.5), Inches(2), Inches(2))
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = RGBColor(0, 82, 184)
        circle1.line.color.rgb = RGBColor(0, 82, 184)
        
        circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(5.5), Inches(2), Inches(2))
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = RGBColor(0, 82, 184)
        circle2.line.color.rgb = RGBColor(0, 82, 184)
        
        # Titre principal - Bien visible et centré
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
        tf = title_box.text_frame
        tf.text = main_title
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Sous-titre
        subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(7), Inches(0.6))
        tf2 = subtitle_box.text_frame
        tf2.text = subtitle
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
        
        # Informations étudiant avec encadré
        info_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.5), Inches(6), Inches(1.3))
        info_bg.fill.solid()
        info_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        info_bg.line.color.rgb = RGBColor(255, 255, 255)
        
        info_box = slide.shapes.add_textbox(Inches(2.2), Inches(5.6), Inches(5.6), Inches(1.1))
        tf3 = info_box.text_frame
        tf3.text = student_info
        tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in tf3.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = COLOR_ACCENT
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.space_after = Pt(6)
        
        return slide
    
    def add_content_slide(title, content_list, slide_number=None):
        """Slide de contenu avec design moderne"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Fond blanc
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(255, 255, 255)
        
        # Bande latérale décorative
        side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = COLOR_PRIMARY
        side_bar.line.color.rgb = COLOR_PRIMARY
        
        # En-tête avec forme
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), 0, Inches(9.7), Inches(1.2))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = COLOR_LIGHT
        header_shape.line.color.rgb = COLOR_LIGHT
        
        # Accent coloré sous le titre
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.1), Inches(9.7), Inches(0.1))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLOR_SECONDARY
        accent_line.line.color.rgb = COLOR_SECONDARY
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(8.5), Inches(0.7))
        tf = title_box.text_frame
        tf.text = title
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        # Numéro de slide si fourni
        if slide_number:
            num_box = slide.shapes.add_textbox(Inches(9.2), Inches(7.1), Inches(0.6), Inches(0.3))
            tf_num = num_box.text_frame
            tf_num.text = str(slide_number)
            p_num = tf_num.paragraphs[0]
            p_num.font.size = Pt(14)
            p_num.font.color.rgb = COLOR_PRIMARY
            p_num.alignment = PP_ALIGN.RIGHT
        
        # Zone de contenu avec fond légèrement coloré
        content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = COLOR_LIGHT
        content_bg.line.color.rgb = COLOR_LIGHT
        
        # Contenu texte
        content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(7.7), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        
        for i, text in enumerate(content_list):
            if i > 0:
                tf.add_paragraph()
            p = tf.paragraphs[i]
            
            # Ajouter des puces pour les items non vides
            if text.strip():
                p.text = "• " + text
                p.font.size = Pt(18)
                p.font.color.rgb = COLOR_ACCENT
                p.space_after = Pt(12)
                p.line_spacing = 1.2
            else:
                p.text = ""
                p.space_after = Pt(6)
        
        return slide
    
    def add_highlight_slide(title, highlight_text, details_list):
        """Slide de mise en avant"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Fond blanc
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(255, 255, 255)
        
        # Bande latérale
        side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = COLOR_PRIMARY
        side_bar.line.color.rgb = COLOR_PRIMARY
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(0.7))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Encadré principal avec ombre
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(2.3), Inches(6.8), Inches(1.8))
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
        shadow.line.color.rgb = RGBColor(200, 200, 200)
        
        highlight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(6.8), Inches(1.8))
        highlight_box.fill.solid()
        highlight_box.fill.fore_color.rgb = COLOR_SECONDARY
        highlight_box.line.color.rgb = COLOR_SECONDARY
        
        # Texte dans l'encadré
        text_box = slide.shapes.add_textbox(Inches(1.8), Inches(2.5), Inches(6.2), Inches(1.2))
        tf2 = text_box.text_frame
        tf2.text = highlight_text
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
        
        # Détails en bas
        details_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(2))
        tf3 = details_box.text_frame
        tf3.word_wrap = True
        
        for i, detail in enumerate(details_list):
            if i > 0:
                tf3.add_paragraph()
            p3 = tf3.paragraphs[i]
            p3.text = "• " + detail
            p3.font.size = Pt(18)
            p3.font.color.rgb = COLOR_ACCENT
            p3.space_after = Pt(10)
            p3.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # ==========================================
    # SLIDE 1 : PAGE DE GARDE
    # ==========================================
    add_title_slide(
        "Retour d'Expérience\nProfessionnelle",
        "Stage Community Manager",
        "Mouhamadou SOW\nB3 Communication Digitale et Internationale\nBuroxia - Février à Avril 2026"
    )
    
    # ==========================================
    # SLIDE 2 : PLAN
    # ==========================================
    add_content_slide(
        "Plan de la Présentation",
        [
            "Introduction : contexte et objectifs du stage",
            "",
            "Présentation de l'entreprise Buroxia",
            "",
            "Mes missions de Community Manager",
            "",
            "Mission phare : collaboration avec l'influenceuse Ayaa.lbns",
            "",
            "Bilan personnel et perspectives professionnelles"
        ],
        2
    )
    
    # ==========================================
    # SLIDE 3 : INTRODUCTION
    # ==========================================
    add_content_slide(
        "Introduction",
        [
            "Stage de 3 mois effectué de février à avril 2026",
            "",
            "Entreprise : Buroxia, centre d'affaires situé à Marseille",
            "",
            "Poste : Community Manager",
            "",
            "Objectif principal : développer la visibilité digitale et la notoriété de Buroxia",
            "",
            "Mission phare : signature d'un partenariat avec l'influenceuse Ayaa.lbns"
        ],
        3
    )
    
    # ==========================================
    # SLIDE 4 : BUROXIA
    # ==========================================
    add_content_slide(
        "Présentation de Buroxia",
        [
            "Centre d'affaires implanté au 24 avenue du Prado, Marseille 6ème",
            "",
            "330 m² d'espaces professionnels avec plus de 15 ans d'expérience",
            "",
            "Services : domiciliation d'entreprises, bureaux équipés, salles de réunion",
            "",
            "Clientèle : entrepreneurs, TPE, professions libérales et influenceurs",
            "",
            "Positionnement : adresse prestigieuse et relation de proximité"
        ],
        4
    )
    
    # ==========================================
    # SLIDE 5 : ENJEUX
    # ==========================================
    add_content_slide(
        "Enjeux et Contexte du Stage",
        [
            "Marché marseillais des centres d'affaires très concurrentiel",
            "",
            "Problématique : structure peu connue malgré une offre de qualité",
            "",
            "Enjeu stratégique : développer rapidement la visibilité et la notoriété",
            "",
            "Solution : stratégie digitale couplée au marketing d'influence",
            "",
            "Tuteur : Stephanne Heyoppe, propriétaire et directeur de Buroxia"
        ],
        5
    )
    
    # ==========================================
    # SLIDE 6 : MISSIONS
    # ==========================================
    add_content_slide(
        "Mes Missions de Community Manager",
        [
            "Gestion des réseaux sociaux : Instagram, TikTok, LinkedIn et Facebook",
            "",
            "Création de contenus visuels et vidéo avec Canva et outils de montage",
            "",
            "Refonte des supports commerciaux : plaquette, grille tarifaire, fiches",
            "",
            "Mise à jour et optimisation du site web WordPress",
            "",
            "Prospection commerciale, rédaction de devis et accompagnement clients"
        ],
        6
    )
    
    # ==========================================
    # SLIDE 7 : MISSION PHARE
    # ==========================================
    add_highlight_slide(
        "Mission Phare",
        "Collaboration avec l'influenceuse Ayaa.lbns",
        [
            "Projet stratégique mené de A à Z",
            "De l'idée initiale à la signature du contrat",
            "Responsabilité : prospection, négociation, coordination"
        ]
    )
    
    # ==========================================
    # SLIDE 8 : GENÈSE
    # ==========================================
    add_content_slide(
        "Genèse du Projet d'Influence",
        [
            "Initiative personnelle présentée lors d'une réunion avec la direction",
            "",
            "Problématique : comment accroître rapidement la visibilité de Buroxia ?",
            "",
            "Stratégie : utiliser le marketing d'influence pour toucher une audience qualifiée",
            "",
            "Validation du projet par la direction qui me confie la mission complète",
            "",
            "Objectif : identifier, contacter et conclure un partenariat avec une influenceuse"
        ],
        8
    )
    
    # ==========================================
    # SLIDE 9 : PROSPECTION
    # ==========================================
    add_content_slide(
        "Prospection d'Influenceurs",
        [
            "Mise en place d'une veille active quotidienne sur Instagram et TikTok",
            "",
            "Analyse de chaque profil : audience, ligne éditoriale, valeurs, adéquation",
            "",
            "Prises de contact par messages privés et e-mails professionnels",
            "",
            "Difficultés : plusieurs refus liés au manque de notoriété de Buroxia",
            "",
            "Persévérance maintenue jusqu'à l'opportunité avec Ayaa.lbns"
        ],
        9
    )
    
    # ==========================================
    # SLIDE 10 : PROFIL AYAA
    # ==========================================
    add_content_slide(
        "Profil de l'Influenceuse Ayaa.lbns",
        [
            "Nom : Ouweik Chrystelle, pseudonyme Ayaa.lbns",
            "",
            "Engagement social fort : aide aux personnes sans abri en France et international",
            "",
            "Actions concrètes : rénovation d'habitats, distribution de nourriture",
            "",
            "Visibilité accrue durant le Ramadan avec ruptures du jeûne collectives",
            "",
            "Image authentique et humaine, alignée avec les valeurs de proximité de Buroxia"
        ],
        10
    )
    
    # ==========================================
    # SLIDE 11 : OPPORTUNITÉ
    # ==========================================
    add_content_slide(
        "L'Opportunité Saisie",
        [
            "Repérage d'une story Instagram : Ayaa.lbns recherche un bureau à Marseille",
            "",
            "Correspondance parfaite avec l'offre de Buroxia",
            "",
            "Contact immédiat du manager via l'adresse e-mail indiquée",
            "",
            "Timing idéal pour proposer une collaboration structurée"
        ],
        11
    )
    
    # ==========================================
    # SLIDE 12 : PREMIER CONTACT
    # ==========================================
    add_content_slide(
        "Premier Contact Professionnel",
        [
            "Rédaction d'un e-mail professionnel au manager d'Ayaa.lbns",
            "",
            "Présentation détaillée de Buroxia et de ses services",
            "",
            "Mise en avant des solutions : domiciliation, bureaux, coworking",
            "",
            "Ouverture vers un partenariat incluant de la visibilité pour Buroxia",
            "",
            "Résultat : réponse favorable et ouverture d'échanges approfondis"
        ],
        12
    )
    
    # ==========================================
    # SLIDE 13 : OUTILS
    # ==========================================
    add_content_slide(
        "Élaboration des Outils Commerciaux",
        [
            "Création d'une fiche de prestations de services complète",
            "",
            "Conception d'une grille tarifaire claire avec formules et promotions",
            "",
            "Design professionnel et cohérent réalisé avec Canva",
            "",
            "Utilisation pour la négociation et les présentations lors des réunions",
            "",
            "Réutilisation pour l'ensemble de la prospection commerciale de Buroxia"
        ],
        13
    )
    
    # ==========================================
    # SLIDE 14 : RÉUNIONS
    # ==========================================
    add_content_slide(
        "Organisation des Réunions",
        [
            "Trois réunions entre l'équipe Buroxia et l'équipe d'Ayaa.lbns",
            "",
            "Réunion 1 (visio) : présentation de l'offre et premiers échanges",
            "",
            "Réunion 2 (visio) : affinage de la proposition et discussions",
            "",
            "Réunion 3 (présentiel) : visite des locaux et finalisation du partenariat",
            "",
            "Mon rôle : préparation des supports, participation active et suivi complet"
        ],
        14
    )
    
    # ==========================================
    # SLIDE 15 : MODÈLE
    # ==========================================
    add_content_slide(
        "Modèle de Collaboration Innovant",
        [
            "Clause spécifique : 50% location + 50% contrepartie publicitaire",
            "",
            "Pour Buroxia : visibilité auprès d'une audience qualifiée",
            "",
            "Pour Ayaa.lbns : espace professionnel à conditions avantageuses",
            "",
            "Modèle gagnant-gagnant ayant permis la signature du contrat"
        ],
        15
    )
    
    # ==========================================
    # SLIDE 16 : RÉSULTATS
    # ==========================================
    add_content_slide(
        "Résultats Obtenus",
        [
            "Signature effective d'une collaboration avec Ayaa.lbns",
            "",
            "Exposition de Buroxia auprès de l'audience engagée de l'influenceuse",
            "",
            "Association à une créatrice reconnue pour son engagement social",
            "",
            "Livrables : fiche de prestations, grille tarifaire, supports, contrat",
            "",
            "Premier partenariat d'influence structuré ouvrant la voie à d'autres"
        ],
        16
    )
    
    # ==========================================
    # SLIDE 17 : VALEUR AJOUTÉE
    # ==========================================
    add_content_slide(
        "Valeur Ajoutée pour l'Entreprise",
        [
            "Concrétisation du premier partenariat influenceur pour Buroxia",
            "",
            "Renforcement de la crédibilité de la marque",
            "",
            "Élargissement de la visibilité sur une cible qualifiée",
            "",
            "Création d'outils commerciaux réutilisables",
            "",
            "Ouverture stratégique vers de futures collaborations d'influence"
        ],
        17
    )
    
    # ==========================================
    # SLIDE 18 : ANALYSE
    # ==========================================
    add_content_slide(
        "Analyse Critique et Axes d'Amélioration",
        [
            "Points forts : initiative personnelle, persévérance, saisie de l'opportunité",
            "",
            "Difficultés : refus initiaux, faible notoriété, coordination à distance",
            "",
            "Axes d'amélioration : systématiser le processus de prospection",
            "",
            "Recommandation : mettre en place des indicateurs de mesure des retombées"
        ],
        18
    )
    
    # ==========================================
    # SLIDE 19 : COMPÉTENCES
    # ==========================================
    add_content_slide(
        "Compétences Développées",
        [
            "Prospection de partenaires influenceurs par veille active et ciblage",
            "",
            "Rédaction de supports commerciaux structurés et professionnels",
            "",
            "Création de contenus visuels et vidéo pour réseaux sociaux",
            "",
            "Organisation et animation de réunions à distance et en présentiel",
            "",
            "Négociation commerciale et construction d'offres gagnant-gagnant"
        ],
        19
    )
    
    # ==========================================
    # SLIDE 20 : SAVOIR-ÊTRE
    # ==========================================
    add_content_slide(
        "Savoir-Être Développé",
        [
            "Persévérance et résilience face aux refus et aux obstacles",
            "",
            "Initiative et capacité à proposer des projets stratégiques",
            "",
            "Prise de responsabilité sur un projet de bout en bout",
            "",
            "Rigueur rédactionnelle et professionnalisme dans les échanges",
            "",
            "Réactivité et saisie d'opportunités en temps réel"
        ],
        20
    )
    
    # ==========================================
    # SLIDE 21 : BILAN
    # ==========================================
    add_content_slide(
        "Bilan Personnel",
        [
            "Expérience professionnelle extrêmement formatrice et concrète",
            "",
            "Confirmation de mon intérêt pour la communication digitale",
            "",
            "Développement de compétences opérationnelles et stratégiques",
            "",
            "Capacité à travailler en autonomie avec collaboration directionnelle",
            "",
            "Réussite d'un projet stratégique de l'idée à la signature"
        ],
        21
    )
    
    # ==========================================
    # SLIDE 22 : PROJET PRO
    # ==========================================
    add_content_slide(
        "Projet Professionnel",
        [
            "Orientation confirmée : communication digitale et marketing d'influence",
            "",
            "Objectif : concevoir et piloter des stratégies de visibilité pour les marques",
            "",
            "Ambition : combiner créativité et performance commerciale",
            "",
            "Souhait d'accompagner les marques dans leur transformation digitale"
        ],
        22
    )
    
    # ==========================================
    # SLIDE 23 : CONCLUSION
    # ==========================================
    add_content_slide(
        "Conclusion",
        [
            "Mission phare réussie : de l'idée initiale à la signature du contrat",
            "",
            "Responsabilité assumée sur un projet stratégique rare pour un stagiaire",
            "",
            "Contribution concrète et mesurable au développement de Buroxia",
            "",
            "Compétences développées en prospection, négociation et création",
            "",
            "Remerciements à Stephanne Heyoppe et à toute l'équipe Buroxia"
        ],
        23
    )
    
    # ==========================================
    # SLIDE 24 : QUESTIONS
    # ==========================================
    slide24 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond bleu
    bg24 = slide24.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg24.fill.solid()
    bg24.fill.fore_color.rgb = COLOR_PRIMARY
    bg24.line.color.rgb = COLOR_PRIMARY
    
    # Cercles décoratifs
    circle24_1 = slide24.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(0.5), Inches(2), Inches(2))
    circle24_1.fill.solid()
    circle24_1.fill.fore_color.rgb = RGBColor(0, 82, 184)
    circle24_1.line.color.rgb = RGBColor(0, 82, 184)
    
    circle24_2 = slide24.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(5.5), Inches(2), Inches(2))
    circle24_2.fill.solid()
    circle24_2.fill.fore_color.rgb = RGBColor(0, 82, 184)
    circle24_2.line.color.rgb = RGBColor(0, 82, 184)
    
    # Texte principal
    main_text24 = slide24.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf24 = main_text24.text_frame
    tf24.text = "Merci de votre attention"
    tf24.vertical_anchor = MSO_ANCHOR.MIDDLE
    p24 = tf24.paragraphs[0]
    p24.font.size = Pt(48)
    p24.font.bold = True
    p24.font.color.rgb = RGBColor(255, 255, 255)
    p24.alignment = PP_ALIGN.CENTER
    
    # Sous-texte
    sub_text24 = slide24.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(0.8))
    tf24b = sub_text24.text_frame
    tf24b.text = "Je suis à votre disposition pour vos questions"
    p24b = tf24b.paragraphs[0]
    p24b.font.size = Pt(28)
    p24b.font.color.rgb = RGBColor(255, 255, 255)
    p24b.alignment = PP_ALIGN.CENTER
    
    # Sauvegarde
    prs.save('/workspace/Soutenance_Mouhamadou_Sow_V2.pptx')
    print("✅ Présentation V2 créée avec succès !")
    print(f"📊 Nombre de slides : {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
